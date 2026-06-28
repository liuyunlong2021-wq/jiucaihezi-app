"""Generate product overview PPT for 韭菜盒子 Studio."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
OLIVE = RGBColor(0x6B, 0x8E, 0x23)
OLIVE_DARK = RGBColor(0x55, 0x6B, 0x2F)
INK = RGBColor(0x2D, 0x2D, 0x2D)
INK2 = RGBColor(0x55, 0x55, 0x55)
INK3 = RGBColor(0x88, 0x88, 0x88)
PAPER = RGBColor(0xEE, 0xE1, 0xCE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD5, 0xC7, 0x87)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def S():
    return prs.slides.add_slide(prs.slide_layouts[6])

def TB(slide, l, t, w, h, text='', fs=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Inter'
    p.alignment = align
    return tf

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def HDR(slide, tag, title, n):
    R(slide, 0, 0, 13.333, 0.06, OLIVE)
    TB(slide, 0.8, 0.5, 11.5, 0.5, tag, fs=13, color=OLIVE, bold=True)
    TB(slide, 0.8, 0.9, 11.5, 1.0, title, fs=32, bold=True, color=INK)
    R(slide, 0.8, 1.8, 2.0, 0.04, OLIVE)
    TB(slide, 12.3, 7.0, 0.8, 0.4, str(n), fs=10, color=INK3, align=PP_ALIGN.RIGHT)

def STRIP(slide):
    R(slide, 0, 7.15, 13.333, 0.35, OLIVE_DARK)
    TB(slide, 0.8, 7.17, 11.5, 0.3, '\u97ed\u83dc\u76d2\u5b50 Studio  |  jiucaihezi.studio', fs=9, color=WHITE)

def CARD(slide, l, t, w, h, icon, title, desc, ac=None):
    if ac is None: ac = OLIVE
    R(slide, l, t, w, h, WHITE)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l+0.3), Inches(t+0.25), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = ac; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = icon; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    TB(slide, l+1.0, t+0.2, w-1.3, 0.4, title, fs=16, bold=True, color=INK)
    TB(slide, l+1.0, t+0.6, w-1.3, h-0.8, desc, fs=12, color=INK2)

def NS(n):
    s = S(); STRIP(s); return s


# ==== SLIDE 1 - COVER ====
s = S()
R(s, 0, 0, 13.333, 7.5, DARK_BG)
R(s, 0, 3.1, 13.333, 0.06, OLIVE)
TB(s, 1.5, 1.8, 10.3, 1.2, '\u97ed\u83dc\u76d2\u5b50 Studio', fs=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
TB(s, 1.5, 3.5, 10.3, 0.8, '\u672c\u5730\u4f18\u5148 \u00b7 \u7eaf\u624b\u52a8 \u00b7 \u5168\u6a21\u578b AI \u5de5\u4f5c\u53f0', fs=22, color=GOLD, align=PP_ALIGN.CENTER)
TB(s, 2.0, 4.6, 9.3, 1.0, '\u4e00\u952e\u767b\u5f55\uff0c\u5373\u523b\u4f7f\u7528 Claude / GPT / Grok \u7b49\u5168\u90e8\u9876\u5c16\u6a21\u578b\n\u81ea\u7531\u7ec4\u5408 Skill \u00b7 \u77e5\u8bc6\u5e93 \u00b7 \u521b\u4f5c\u5de5\u5177', fs=15, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)
TB(s, 1.5, 6.2, 10.3, 0.4, 'v1.0.15  |  jiucaihezi.studio', fs=12, color=INK3, align=PP_ALIGN.CENTER)

# ==== SLIDE 2 - POSITIONING ====
s = NS(2)
HDR(s, '\u4ea7\u54c1\u5b9a\u4f4d', 'AI \u5de5\u4f5c\u53f0\uff0c\u4e0d\u662f AI Agent', 2)
TB(s, 0.8, 2.3, 11.5, 0.6, '\u300c\u7528\u6237\u4e3b\u5bfc\uff0cAI \u6267\u884c\u300d\u2014\u2014 \u7528\u6237\u9009\u62e9\u6a21\u578b\u3001\u9009\u62e9 Skill\u3001\u9009\u62e9\u5de5\u5177\uff0cAI \u6309\u663e\u5f0f\u914d\u7f6e\u7cbe\u51c6\u6267\u884c\u3002\u4e0d\u641e\u9ed1\u76d2\u81ea\u52a8\u5316\u3002', fs=15, color=INK2)
cards = [
    ('\U0001f464', '\u666e\u901a\u7528\u6237', '\u6253\u5f00\u5373\u7528\uff0c\u4e00\u952e\u767b\u5f55\u5c31\u80fd\u548c\n\u9876\u5c16\u6a21\u578b\u5bf9\u8bdd\u3001\u751f\u56fe\u3001\u505a\u89c6\u9891'),
    ('\u270d\ufe0f', '\u5185\u5bb9\u521b\u4f5c\u8005', 'Skill \u7cfb\u7edf\u8c03\u7528\u4e13\u4e1a\u80fd\u529b\u5305\n\u642d\u77e5\u8bc6\u5e93\u9632 AI \u5931\u5fc6'),
    ('\U0001f4bb', '\u5f00\u53d1\u8005', '\u672c\u5730\u9879\u76ee\u8bfb\u5199\u3001\u547d\u4ee4\u6267\u884c\nOpenCode \u6587/\u6b66\u53cc\u6a21\u5f0f'),
    ('\U0001f3a8', '\u8bbe\u8ba1\u5e08', '\u521b\u4f5c\u9762\u677f\u4e00\u7ad9\u5f0f\u751f\u56fe\u3001\n\u751f\u89c6\u9891\u3001\u751f\u97f3\u9891'),
]
for i, (icon, title, desc) in enumerate(cards):
    CARD(s, 0.8 + i*3.05, 3.2, 2.85, 2.0, icon, title, desc)
TB(s, 0.8, 5.8, 11.5, 0.8, '\u6838\u5fc3\u539f\u5219\uff1aSkill \u51b3\u5b9a AI \u662f\u8c01  \u00b7  \u77e5\u8bc6\u5e93\u63d0\u4f9b\u53c2\u8003\u8d44\u6599  \u00b7  \u5de5\u5177\u4ed3\u5e93\u63d0\u4f9b\u672c\u5730\u624b\u811a', fs=13, color=INK3, align=PP_ALIGN.CENTER)

# ==== SLIDE 3 - ONE CLICK ====
s = NS(3)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2460', '\u96f6\u95e8\u69db\uff0c\u4e00\u952e\u767b\u5f55', 3)
for i, item in enumerate(['\u4e0b\u8f7d\u5e94\u7528', '\u4e00\u952e\u767b\u5f55', '\u5373\u523b\u4f7f\u7528\u5168\u90e8\u9876\u5c16\u6a21\u578b']):
    x = 2.0 + i*3.5
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.5), Inches(2.6), Inches(1.2), Inches(1.2))
    c.fill.solid(); c.fill.fore_color.rgb = OLIVE if i != 1 else GOLD; c.line.fill.background()
    tf = c.text_frame; p = tf.paragraphs[0]; p.text = item; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    if i < 2:
        TB(s, x+1.9, 2.9, 1.4, 0.5, '\u2192', fs=28, color=INK3, align=PP_ALIGN.CENTER)
bullets = [
    ('\U0001f511 \u4e0d\u9700\u8981 API Key', '\u8d26\u53f7\u767b\u5f55\u5373\u7528\uff0c\u5185\u7f6e NewAPI \u4e2d\u8f6c\uff0c\u7edf\u4e00\u9274\u6743'),
    ('\U0001f310 \u4e0d\u9700\u8981\u7ffb\u5899', '\u672c\u5730\u684c\u9762\u5e94\u7528\uff0c\u76f4\u8fde\u56fd\u5185\u52a0\u901f\u8282\u70b9'),
    ('\u2699\ufe0f \u4e0d\u9700\u8981\u914d\u7f6e', '\u6a21\u578b\u5217\u8868\u81ea\u52a8\u540c\u6b65\uff0c\u4e00\u952e\u5207\u6362 30+ \u6a21\u578b'),
    ('\U0001f4cb \u4e00\u952e\u6284\u914d\u7f6e', '\u767b\u5f55\u540e\u590d\u5236\u5168\u5957 API \u4fe1\u606f\uff0c\u5728\u4efb\u4f55 AI \u5ba2\u6237\u7aef\u4f7f\u7528'),
]
for i, (title, desc) in enumerate(bullets):
    y = 4.2 + i*0.7
    TB(s, 2.0, y, 4.0, 0.35, title, fs=14, bold=True, color=INK)
    TB(s, 6.0, y, 6.0, 0.35, desc, fs=13, color=INK2)

# ==== SLIDE 4 - MODELS ====
s = NS(4)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2461', '\u5168\u6a21\u578b\u81ea\u7531\u5207\u6362', 4)
models = [
    ('\U0001f4ac \u6587\u672c\u5bf9\u8bdd', 'Claude Sonnet 4.6\u3001GPT-4o\u3001Grok-4\u3001Gemini 2.5 Pro\u3001DeepSeek V3'),
    ('\U0001f3a8 \u56fe\u7247\u751f\u6210', 'GPT Image 2\u3001Nano Banana\u3001FLUX Klein 9B'),
    ('\U0001f3ac \u89c6\u9891\u751f\u6210', 'Veo 3\u3001Grok Video\u3001Seedance 2.0'),
    ('\U0001f3b5 \u97f3\u9891\u751f\u6210', 'Suno V5\u3001RH \u58f0\u97f3'),
]
for i, (cat, mods) in enumerate(models):
    y = 2.3 + i*1.1
    R(s, 1.0, y, 11.3, 0.85, WHITE)
    TB(s, 1.2, y+0.1, 2.5, 0.35, cat, fs=16, bold=True, color=OLIVE_DARK)
    TB(s, 3.8, y+0.1, 8.3, 0.6, mods, fs=14, color=INK)
TB(s, 1.0, 6.2, 5.5, 0.4, '\U0001f9e0 \u6a21\u578b\u611f\u77e5\uff1a\u81ea\u52a8\u68c0\u6d4b vision / text-only\uff0c\u56fe\u7247\u81ea\u9002\u5e94\u683c\u5f0f', fs=12, color=INK2)
TB(s, 7.0, 6.2, 5.5, 0.4, '\U0001f4ca Token \u6c34\u4f4d\u8ba1\uff1a\u5b9e\u65f6\u663e\u793a\u4e0a\u4e0b\u6587\u7528\u91cf\uff0c\u4e09\u8272\u9884\u8b66', fs=12, color=INK2)

# ==== SLIDE 5 - SKILL ====
s = NS(5)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2462', 'Skill \u7cfb\u7edf \u2014 \u7ed9 AI \u88c5\u4e0a\u4e13\u4e1a\u5927\u8111', 5)
TB(s, 0.8, 2.2, 6.0, 0.8, 'Skill \u662f\u5b98\u65b9 Anthropic Skill \u6807\u51c6\u7684\u80fd\u529b\u5305\u3002\n\u6bcf\u4e2a Skill = SKILL.md + \u53c2\u8003\u8d44\u6599 + \u811a\u672c\uff0c\n\u8ba9 AI \u77ac\u95f4\u83b7\u5f97\u4e13\u4e1a\u9886\u57df\u80fd\u529b\u3002', fs=14, color=INK2)
skill_list = [
    ['\u2696\ufe0f \u5f8b\u5e08\u5de5\u4f5c\u53f0', '\U0001f4dd \u6f2b\u5267\u5267\u672c\u751f\u6210\u5668', '\U0001f3a8 \u54c1\u724c\u6307\u5357'],
    ['\U0001f4bb Claude API', '\U0001f3d7\ufe0f MCP \u6784\u5efa\u5668', '\U0001f58c\ufe0f \u524d\u7aef\u8bbe\u8ba1'],
    ['\U0001f4c4 Word/PDF/PPT', '\U0001f4ca Excel', '\U0001f91d \u6587\u6863\u534f\u4f5c'],
    ['\U0001f4da \u77e5\u8bc6\u5e93\u67b6\u6784\u5e08', '\U0001f527 Skill \u521b\u5efa\u5668', '\U0001f4ee \u5185\u90e8\u901a\u8baf'],
]
for ri, row in enumerate(skill_list):
    for ci, skill in enumerate(row):
        x, y = 0.8 + ci*2.1, 3.3 + ri*0.9
        R(s, x, y, 1.95, 0.7, WHITE)
        TB(s, x+0.1, y+0.1, 1.75, 0.5, skill, fs=12, color=INK, align=PP_ALIGN.CENTER)
TB(s, 7.5, 2.2, 5.0, 3.5, '\u2726 \u4e2d\u592e Skill \u5e93\uff1a20+ \u5b98\u65b9 Skill\n    GitHub \u4e00\u952e\u5bfc\u5165\u793e\u533a Skill\n\n\u2726 \u7528\u6237\u81ea\u5efa\uff1a\u4ece\u5bf9\u8bdd\u3001\u7f16\u8f91\u5668\u3001\n    \u77e5\u8bc6\u5e93\u63d0\u70bc\u7ecf\u9a8c \u2192 \u751f\u6210 Skill\n\n\u2726 \u7eaf\u624b\u52a8\uff1a\u7528\u6237\u9009\u62e9 Skill \u2192\n    AI \u6309 SKILL.md \u6267\u884c', fs=13, color=INK2)

# ==== SLIDE 6 - KNOWLEDGE BASE ====
s = NS(6)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2463', '\u77e5\u8bc6\u5e93 \u2014 \u7ed9 AI \u88c5\u4e0a\u957f\u671f\u8bb0\u5fc6', 6)
TB(s, 0.8, 2.1, 11.5, 0.5, '\U0001f4a1 \u89e3\u51b3 AI\u300c\u804a\u5b8c\u5c31\u5fd8\u300d\u7684\u6838\u5fc3\u75db\u70b9 \u2014\u2014 \u7528\u6237\u4e0a\u4f20\u8d44\u6599 \u2192 \u7cfb\u7edf\u6574\u7406\u4e3a\u7ed3\u6784\u5316\u77e5\u8bc6\u9875 \u2192 \u5bf9\u8bdd\u65f6\u81ea\u52a8\u53ec\u56de', fs=14, color=INK2)
layers = [
    ('raw/', '\u539f\u59cb\u7d20\u6750', '\u7528\u6237\u62d6\u62fd\u4e0a\u4f20\n\u4e0d\u6539\u52a8\u539f\u6587\u4ef6', OLIVE),
    ('wiki/', '\u7ed3\u6784\u5316\u77e5\u8bc6\u9875', '\u53ef\u68c0\u7d22\u3001\u53ef\u5f15\u7528\nAI \u7eed\u5199\u4e0d\u5931\u5fc6', GOLD),
    ('CLAUDE.md', '\u8bb0\u5fc6\u951a\u70b9', '\u5b9a\u4e49\u77e5\u8bc6\u5e93\u89c4\u5219\n\u77e5\u8bc6\u5e93\u7684\u300c\u8bf4\u660e\u4e66\u300d', OLIVE_DARK),
]
for i, (name, title, desc, color) in enumerate(layers):
    x = 1.2 + i*4.0
    R(s, x, 3.0, 3.5, 2.2, WHITE)
    R(s, x, 3.0, 3.5, 0.5, color)
    TB(s, x+0.1, 3.05, 3.3, 0.4, f'{name}  {title}', fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    TB(s, x+0.2, 3.7, 3.1, 1.2, desc, fs=13, color=INK2, align=PP_ALIGN.CENTER)
TB(s, 1.0, 5.6, 5.5, 0.7, '\U0001f6e1\ufe0f \u5b89\u5168\u5e95\u7ebf\uff1a\u77e5\u8bc6\u5e93\u53ea\u63a5\u53d7\u7528\u6237\u624b\u52a8\u6dfb\u52a0\n    AI \u7981\u6b62\u81ea\u52a8\u5199\u5165\uff0c\u675c\u7edd\u5e7b\u89c9\u6c61\u67d3', fs=12, color=INK2)
TB(s, 7.0, 5.6, 5.5, 0.7, '\U0001f4e6 \u5f00\u7bb1\u6a21\u677f\uff1a\u5f8b\u5e08\u6848\u4ef6\u5e93 / \u5c0f\u8bf4\u77e5\u8bc6\u5e93\n    \u6f2b\u5267\u5267\u672c\u9879\u76ee\uff0c\u4e00\u952e\u521b\u5efa', fs=12, color=INK2)

# ==== SLIDE 7 - CREATION ====
s = NS(7)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2464', '\u521b\u4f5c\u9762\u677f \u2014 \u4e00\u7ad9\u5f0f AIGC \u5de5\u5382', 7)
gen = [
    ('\U0001f5bc\ufe0f \u56fe\u7247\u751f\u6210', '\u6587\u751f\u56fe / \u56fe\u751f\u56fe\n\u6bd4\u4f8b \u00b7 \u5206\u8fa8\u7387 \u00b7 \u98ce\u683c\u81ea\u7531\u63a7\u5236\nGPT Image 2 / Nano Banana / FLUX', OLIVE),
    ('\U0001f3ac \u89c6\u9891\u751f\u6210', 'Veo 3 / Grok Video / Seedance 2.0\n\u4e09\u5927\u5f15\u64ce\u968f\u610f\u5207\u6362\n\u6bd4\u4f8b \u00b7 \u65f6\u957f\u53ef\u8c03', OLIVE_DARK),
    ('\U0001f3b5 \u97f3\u9891\u751f\u6210', 'Suno V5 \u97f3\u4e50\u751f\u6210\n\u6807\u9898 \u00b7 \u6807\u7b7e \u00b7 MV \u6a21\u5f0f\nRH \u58f0\u97f3', GOLD),
]
for i, (title, desc, color) in enumerate(gen):
    x = 0.8 + i*4.1
    R(s, x, 2.3, 3.8, 2.5, WHITE)
    R(s, x, 2.3, 3.8, 0.06, color)
    TB(s, x+0.3, 2.5, 3.2, 0.5, title, fs=20, bold=True, color=INK)
    TB(s, x+0.3, 3.1, 3.2, 1.5, desc, fs=13, color=INK2)
TB(s, 0.8, 5.2, 11.5, 0.5, '13 \u4e2a\u5a92\u4f53\u6a21\u578b  \u00b7  \u7edf\u4e00\u9274\u6743  \u00b7  \u7edf\u4e00\u5165\u53e3  \u00b7  \u751f\u6210\u7ed3\u679c\u81ea\u52a8\u672c\u5730\u4fdd\u5b58', fs=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
TB(s, 0.8, 5.8, 11.5, 0.5, '\u753b\u5eca\u81ea\u52a8\u5f52\u6863 \u2192 \u4e00\u952e\u4e0b\u8f7d \u2192 \u590d\u5236\u94fe\u63a5 \u2192 \u706f\u7bb1\u9884\u89c8', fs=13, color=INK2, align=PP_ALIGN.CENTER)

# ==== SLIDE 8 - LOCAL TOOLS ====
s = NS(8)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2465', '\u672c\u5730\u5de5\u5177 \u2014 \u684c\u9762\u7aef\u7684\u771f\u6b63\u529b\u91cf', 8)
tools = [
    ('\U0001f4c4', '\u6587\u6863\u8f6c\u6362', 'Office / PDF \u8f6c Markdown\n\u672c\u5730\u5904\u7406\u4e0d\u6cc4\u5bc6'),
    ('\U0001f310', '\u6d4f\u89c8\u5668\u63a7\u5236', '\u53ef\u89c1 Chrome \u641c\u7d22\u3001\u6253\u5f00\u7f51\u9875\n\u622a\u56fe\u3001\u70b9\u51fb\u64cd\u4f5c'),
    ('\U0001f4bb', '\u9879\u76ee\u8bfb\u5199', '\u9879\u76ee\u76ee\u5f55\u5185\u8bfb\u5199\u6587\u4ef6\n\u641c\u7d22\u3001\u6267\u884c\u547d\u4ee4'),
    ('\U0001f3ac', '\u5a92\u4f53\u5904\u7406', 'ffmpeg \u89c6\u9891\u5904\u7406\nyt-dlp \u7f51\u9875\u5a92\u4f53\u91c7\u96c6'),
    ('\U0001f50c', 'OpenCode', '\u6587\u6a21\u5f0f / \u6b66\u6a21\u5f0f\u53cc\u5f15\u64ce\n\u9879\u76ee\u7ea7 AI \u534f\u4f5c'),
]
for i, (icon, title, desc) in enumerate(tools):
    CARD(s, 0.8 + i*2.5, 2.3, 2.3, 2.5, icon, title, desc, OLIVE_DARK if i == 4 else OLIVE)
TB(s, 0.8, 5.3, 11.5, 0.8, '\u684c\u9762\u5e94\u7528\u4e0d\u53ea\u662f\u7f51\u9875\u5957\u58f3\u3002\u97ed\u83dc\u76d2\u5b50\u76f4\u63a5\u8c03\u7528\u64cd\u4f5c\u7cfb\u7edf\u80fd\u529b\uff0c\n\u6240\u6709\u5904\u7406\u5728\u672c\u5730\u5b8c\u6210\uff0c\u6570\u636e\u4e0d\u7ecf\u8fc7\u4e91\u7aef\u3002', fs=13, color=INK2, align=PP_ALIGN.CENTER)

# ==== SLIDE 9 - DUAL + TECH ====
s = NS(9)
HDR(s, '\u6838\u5fc3\u7279\u70b9 \u2466 + \u6280\u672f\u4eae\u70b9', '\u53cc\u7aef\u8986\u76d6 \u00b7 \u6280\u672f\u5b9e\u529b', 9)
R(s, 0.8, 2.2, 5.5, 2.0, WHITE)
R(s, 0.8, 2.2, 5.5, 0.06, OLIVE)
TB(s, 1.0, 2.4, 5.1, 0.4, '\U0001f5a5\ufe0f \u684c\u9762\u7aef', fs=17, bold=True, color=INK)
TB(s, 1.0, 2.9, 5.1, 1.1, 'macOS / Windows \u539f\u751f\u5e94\u7528\n\u5168\u529f\u80fd\uff1a\u672c\u5730\u5de5\u5177 + OpenCode \u6587/\u6b66\u6a21\u5f0f + \u76f4\u8fde\u6a21\u5f0f', fs=13, color=INK2)
R(s, 7.0, 2.2, 5.5, 2.0, WHITE)
R(s, 7.0, 2.2, 5.5, 0.06, OLIVE)
TB(s, 7.2, 2.4, 5.1, 0.4, '\U0001f310 Web \u7aef', fs=17, bold=True, color=INK)
TB(s, 7.2, 2.9, 5.1, 1.1, '\u6d4f\u89c8\u5668\u8bbf\u95ee jiucaihezi.studio\n\u8f7b\u91cf\u76f4\u8fde\u6a21\u5f0f\uff0c\u968f\u65f6\u968f\u5730\u6253\u5f00\u5373\u7528', fs=13, color=INK2)
TB(s, 0.8, 4.0, 11.5, 0.4, '\u53cc\u7aef\u5171\u4eab\u6838\u5fc3\u4f53\u9a8c\uff1a\u6a21\u578b \u00b7 Skill \u00b7 \u521b\u4f5c\u9762\u677f \u00b7 \u77e5\u8bc6\u5e93\uff0c\u89c6\u89c9\u548c\u4ea4\u4e92\u4fdd\u6301\u4e00\u81f4', fs=12, color=INK3, align=PP_ALIGN.CENTER)
tech = [
    ('\u672c\u5730\u4f18\u5148', '\u6570\u636e\u5b58\u672c\u5730 SQLite\nSession Token 0600 \u6743\u9650\u6587\u4ef6'),
    ('\u6d41\u5f0f\u54cd\u5e94', 'SSE \u5b9e\u65f6\u6d41\u5f0f\u8f93\u51fa\n\u9010\u5b57\u663e\u793a\uff0c\u5bf9\u6807 ChatGPT'),
    ('\u5b89\u5168\u67b6\u6784', 'Rust + Tauri \u684c\u9762\u58f3\n\u6743\u9650\u6700\u5c0f\u5316\uff0c\u8def\u5f84\u7a7f\u8d8a\u9632\u62a4'),
    ('\u79bb\u7ebf\u964d\u7ea7', '\u5b58\u50a8\u5c42\u6545\u969c\u81ea\u52a8\u964d\u7ea7\nUI \u4e0d\u767d\u5c4f\uff0c\u6570\u636e\u4e0d\u4e22\u5931'),
    ('\u4e09\u5e73\u53f0\u53d1\u5e03', 'macOS ARM + Intel\nWindows x64 CI \u81ea\u52a8\u6784\u5efa'),
]
for i, (title, desc) in enumerate(tech):
    x = 0.8 + i*2.5
    R(s, x, 4.8, 2.3, 1.8, WHITE)
    TB(s, x+0.15, 4.9, 2.0, 0.35, title, fs=14, bold=True, color=OLIVE_DARK, align=PP_ALIGN.CENTER)
    TB(s, x+0.15, 5.3, 2.0, 1.1, desc, fs=11, color=INK2, align=PP_ALIGN.CENTER)

# ==== SLIDE 10 - COMPARISON ====
s = NS(10)
HDR(s, '\u7ade\u54c1\u5bf9\u6bd4', '\u4e3a\u4ec0\u4e48\u9009\u62e9\u97ed\u83dc\u76d2\u5b50\uff1f', 10)
headers = ['\u5bf9\u6bd4\u7ef4\u5ea6', 'ChatGPT / Claude \u7f51\u9875\u7248', '\u5176\u4ed6 AI \u684c\u9762\u5ba2\u6237\u7aef', '\u97ed\u83dc\u76d2\u5b50 Studio']
col_w = [2.0, 3.3, 3.3, 3.3]
col_x = [1.0]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)
for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    bg = OLIVE_DARK if i == 3 else (OLIVE if i == 0 else INK2)
    R(s, x+0.1, 2.2, w-0.2, 0.45, bg)
    TB(s, x+0.15, 2.22, w-0.3, 0.4, h, fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows = [
    ['\u6a21\u578b\u6570\u91cf', '\u5355\u4e00\u5382\u5546', '\u9700\u624b\u52a8\u914d Key', '\u2705 30+ \u6a21\u578b\u4e00\u952e\u5207\u6362'],
    ['\u4f7f\u7528\u95e8\u69db', '\u9700\u7ffb\u5899 + \u4ed8\u8d39', '\u9700\u81ea\u884c\u7533\u8bf7 API Key', '\u2705 \u4e00\u952e\u767b\u5f55\u5373\u7528'],
    ['\u77e5\u8bc6\u5e93', '\u65e0\u6216\u6709\u9650', '\u901a\u5e38\u65e0', '\u2705 \u4e09\u5c42\u67b6\u6784 + \u9632\u5e7b\u89c9\u6c61\u67d3'],
    ['Skill \u7cfb\u7edf', '\u65e0', '\u65e0', '\u2705 \u5b98\u65b9 Anthropic Skill \u6807\u51c6'],
    ['\u672c\u5730\u5de5\u5177', '\u65e0', '\u6709\u9650', '\u2705 \u6d4f\u89c8\u5668 + \u9879\u76ee\u8bfb\u5199 + \u5a92\u4f53\u5904\u7406'],
    ['\u5a92\u4f53\u751f\u6210', '\u4ec5\u56fe\u7247', '\u9700\u989d\u5916\u914d\u7f6e', '\u2705 \u56fe/\u89c6\u9891/\u97f3\u9891\u4e00\u7ad9\u5f0f'],
    ['\u6570\u636e\u9690\u79c1', '\u5168\u5728\u4e91\u7aef', '\u53d6\u51b3\u4e8e\u5b9e\u73b0', '\u2705 \u672c\u5730 SQLite \u5bc6\u94a5 0600'],
    ['\u53cc\u7aef\u8986\u76d6', '\u4ec5 Web', '\u4ec5\u684c\u9762', '\u2705 \u684c\u9762 + Web \u53cc\u7aef'],
]
for ri, row in enumerate(rows):
    y = 2.75 + ri*0.52
    bg = WHITE if ri % 2 == 0 else PAPER
    R(s, col_x[0]+0.1, y, sum(col_w)-0.2, 0.48, bg)
    for ci, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        c = OLIVE_DARK if ci == 3 else INK
        b = (ci == 3 or ci == 0)
        TB(s, x+0.15, y+0.05, w-0.3, 0.38, cell, fs=11, bold=b, color=c, align=PP_ALIGN.CENTER)

# ==== SLIDE 11 - CLOSING ====
s = S()
R(s, 0, 0, 13.333, 7.5, DARK_BG)
R(s, 0, 3.5, 13.333, 0.06, OLIVE)
TB(s, 1.5, 1.5, 10.3, 1.0, '\u4ea7\u54c1\u73b0\u72b6', fs=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, item in enumerate([
    '\u7248\u672c v1.0.15  \u2014  \u6838\u5fc3\u529f\u80fd\u7a33\u5b9a\uff0c\u8fdb\u5165\u4ea7\u54c1\u6536\u53e3\u671f',
    'macOS (ARM + Intel)  +  Windows x64  \u6b63\u5f0f\u53d1\u5e03',
    'Web \u7aef  jiucaihezi.studio  \u5728\u7ebf\u53ef\u7528',
]):
    TB(s, 3.0, 2.5 + i*0.5, 7.3, 0.4, f'\u2726  {item}', fs=16, color=GOLD, align=PP_ALIGN.CENTER)
TB(s, 1.5, 4.2, 10.3, 1.5, '\u4e0d\u6298\u817e\uff0c\u6253\u5f00\u5c31\u7528\u3002\n\u4f60\u7684\u7b2c\u4e00\u4e2a\u4e5f\u662f\u6700\u540e\u4e00\u4e2a AI \u5de5\u4f5c\u53f0\u3002', fs=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
TB(s, 1.5, 6.0, 10.3, 0.5, '\u97ed\u83dc\u76d2\u5b50 Studio  |  jiucaihezi.studio  |  GitHub: liuyunlong2021-wq/jiucaihezi-app', fs=13, color=INK3, align=PP_ALIGN.CENTER)

# Save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'docs', '\u97ed\u83dc\u76d2\u5b50Studio-\u4ea7\u54c1\u4ecb\u7ecd.pptx')
prs.save(out)
print(f'Done: {out}')
print(f'Slides: {len(prs.slides)}')
